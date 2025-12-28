"""
Comprehensive PowerPoint Content Extraction and Question Generation
Extracts content from 150+ slides and generates bilingual questions
"""

from pptx import Presentation
import os
import json
from datetime import datetime
import re

print("=" * 70)
print("NEPHROPATHOLOGY CONTENT EXTRACTION & QUESTION GENERATION")
print("=" * 70)
print()

# Create output directories
os.makedirs('extracted_content', exist_ok=True)
os.makedirs('extracted_images', exist_ok=True)
os.makedirs('generated_questions', exist_ok=True)

print("Loading PowerPoint presentations...")
prs_en = Presentation('SlidesForSelfStudy_Nephropathology_EN_2024.pptx')
prs_lt = Presentation('SlidesForSelfStudy_Nephropathology_LT_2024.pptx')

print(f"✓ English: {len(prs_en.slides)} slides")
print(f"✓ Lithuanian: {len(prs_lt.slides)} slides")
print()

# ============================================================================
# STAGE 1: EXTRACT ALL SLIDE CONTENT
# ============================================================================

print("STAGE 1: Extracting slide content...")
print("-" * 70)

slide_content = []
total_images = 0

for i, (slide_en, slide_lt) in enumerate(zip(prs_en.slides, prs_lt.slides[:len(prs_en.slides)])):
    # Extract English text
    en_text = []
    for shape in slide_en.shapes:
        if hasattr(shape, 'text') and shape.text.strip():
            en_text.append(shape.text.strip())

    # Extract Lithuanian text
    lt_text = []
    for shape in slide_lt.shapes:
        if hasattr(shape, 'text') and shape.text.strip():
            lt_text.append(shape.text.strip())

    # Count images
    img_count = sum(1 for shape in slide_en.shapes if hasattr(shape, 'image'))
    total_images += img_count

    slide_data = {
        'slide_number': i + 1,
        'en_text': en_text,
        'lt_text': lt_text,
        'image_count': img_count,
        'has_images': img_count > 0
    }

    slide_content.append(slide_data)

    if (i + 1) % 25 == 0:
        print(f"  Processed {i + 1} slides...")

print(f"✓ Extracted content from {len(slide_content)} slides")
print(f"✓ Found {total_images} images total")
print()

# Save slide content
with open('extracted_content/all_slides.json', 'w', encoding='utf-8') as f:
    json.dump(slide_content, f, indent=2, ensure_ascii=False)

print("✓ Saved: extracted_content/all_slides.json")
print()

# ============================================================================
# STAGE 2: IDENTIFY TEACHING TOPICS
# ============================================================================

print("STAGE 2: Identifying teaching topics...")
print("-" * 70)

# Keywords for different diseases/topics
disease_keywords = {
    'MCD': ['minimal change', 'foot process effacement', 'podocyte'],
    'FSGS': ['focal segmental', 'sclerosis', 'fsgs'],
    'MGN': ['membranous', 'spikes', 'subepithelial'],
    'IgAN': ['iga nephropathy', 'mesangial iga', 'henoch'],
    'APIGN': ['postinfectious', 'starry sky', 'humps'],
    'MPGN': ['membranoproliferative', 'double contour', 'tram track'],
    'ANCA_GN': ['anca', 'wegener', 'granulomatosis', 'vasculitis'],
    'ANTI_GBM': ['goodpasture', 'anti-gbm', 'linear'],
    'DIABETIC': ['diabetic', 'kimmelstiel', 'wilson'],
    'AMYLOID': ['amyloid', 'congo red', 'apple green'],
    'ALPORT': ['alport', 'basket weave', 'hearing'],
    'ATN': ['tubular necrosis', 'atn', 'acute tubular'],
    'AIN': ['interstitial nephritis', 'ain'],
    'LUPUS': ['lupus', 'sle', 'wire loop']
}

# Categorize slides by topic
categorized_slides = {disease: [] for disease in disease_keywords.keys()}
categorized_slides['OTHER'] = []

for slide in slide_content:
    all_text = ' '.join(slide['en_text']).lower()
    categorized = False

    for disease, keywords in disease_keywords.items():
        if any(keyword in all_text for keyword in keywords):
            categorized_slides[disease].append(slide['slide_number'])
            categorized = True
            break

    if not categorized and slide['en_text']:  # Has content but not categorized
        categorized_slides['OTHER'].append(slide['slide_number'])

# Print categorization
for disease, slides in categorized_slides.items():
    if slides:
        print(f"  {disease}: {len(slides)} slides")

print()
print("✓ Slides categorized by topic")
print()

# Save categorization
with open('extracted_content/slide_categories.json', 'w', encoding='utf-8') as f:
    json.dump(categorized_slides, f, indent=2)

print("✓ Saved: extracted_content/slide_categories.json")
print()

# ============================================================================
# STAGE 3: EXTRACT KEY FACTS FOR QUESTIONS
# ============================================================================

print("STAGE 3: Extracting key facts...")
print("-" * 70)

key_facts = []

for slide in slide_content:
    if not slide['en_text'] or len(slide['en_text']) < 2:
        continue

    # Look for bullet points or numbered lists
    for i, text_en in enumerate(slide['en_text']):
        # Get corresponding Lithuanian text if available
        text_lt = slide['lt_text'][i] if i < len(slide['lt_text']) else ""

        # Skip titles and very short text
        if len(text_en.split()) < 5:
            continue

        # Look for characteristic patterns
        if any(word in text_en.lower() for word in ['characteristic', 'feature', 'finding', 'shows', 'presents with']):
            key_facts.append({
                'slide': slide['slide_number'],
                'en': text_en,
                'lt': text_lt,
                'has_image': slide['has_images']
            })

print(f"✓ Extracted {len(key_facts)} key facts")
print()

# Save key facts
with open('extracted_content/key_facts.json', 'w', encoding='utf-8') as f:
    json.dump(key_facts, f, indent=2, ensure_ascii=False)

print("✓ Saved: extracted_content/key_facts.json")
print()

# ============================================================================
# SUMMARY
# ============================================================================

print("=" * 70)
print("EXTRACTION COMPLETE!")
print("=" * 70)
print()
print("Summary:")
print(f"  - Slides processed: {len(slide_content)}")
print(f"  - Images found: {total_images}")
print(f"  - Key facts: {len(key_facts)}")
print(f"  - Topics identified: {sum(1 for v in categorized_slides.values() if v)}")
print()
print("Output files:")
print("  ✓ extracted_content/all_slides.json")
print("  ✓ extracted_content/slide_categories.json")
print("  ✓ extracted_content/key_facts.json")
print()
print("Next: Run image extraction and question generation scripts")
print("=" * 70)
